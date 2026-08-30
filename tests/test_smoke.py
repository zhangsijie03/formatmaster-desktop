"""全功能冒烟测试（固化的正式用例）。

覆盖：文档/图片/音频/视频/PDF/识别/实用工具/任务链路/文件安全。
运行：venv/Scripts/python -m pytest tests/test_smoke.py -q
"""
import os
import shutil
import subprocess
import sys
import tempfile

import pytest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

TMP = tempfile.mkdtemp(prefix="fm_smoke_")
FF = None


def _ff():
    global FF
    if FF is None:
        from utils.config import get_ffmpeg_path
        FF = get_ffmpeg_path()
    return FF


def _make_video(path, duration=1, size="160x120", rate=10):
    ff = _ff()
    r = subprocess.run(
        [ff, "-y", "-f", "lavfi", "-i", f"testsrc=duration={duration}:size={size}:rate={rate}",
         "-c:v", "libx264", "-pix_fmt", "yuv420p", path],
        capture_output=True, timeout=60)
    assert r.returncode == 0, r.stderr[-200:]


def _make_png(path, size=(200, 120), color=(200, 100, 50), text=None):
    from PIL import Image, ImageDraw
    img = Image.new("RGB", size, color)
    if text:
        ImageDraw.Draw(img).text((10, size[1] // 2), text, fill="black")
    img.save(path)


@pytest.fixture(scope="module", autouse=True)
def _cleanup():
    yield
    shutil.rmtree(TMP, ignore_errors=True)


class TestDocument:
    def test_docx_txt(self):
        from core.doc_converter import DocumentConverter
        from docx import Document
        src = os.path.join(TMP, "d.docx")
        d = Document(); d.add_paragraph("hello 测试"); d.save(src)
        out = os.path.join(TMP, "d.txt")
        assert DocumentConverter().convert(src, out)
        assert os.path.isfile(out)

    def test_xlsx_csv(self):
        from core.doc_converter import DocumentConverter
        from openpyxl import Workbook
        src = os.path.join(TMP, "x.xlsx")
        wb = Workbook(); wb.active["A1"] = "v"; wb.save(src)
        out = os.path.join(TMP, "x.csv")
        assert DocumentConverter().convert(src, out)
        assert os.path.isfile(out)

    def test_pdf_txt(self):
        from core.doc_converter import DocumentConverter
        from reportlab.pdfgen import canvas
        src = os.path.join(TMP, "p.pdf")
        c = canvas.Canvas(src); c.drawString(50, 700, "abc"); c.save()
        out = os.path.join(TMP, "p.txt")
        assert DocumentConverter().convert(src, out)

    def test_pptx_txt(self):
        from core.doc_converter import DocumentConverter
        from pptx import Presentation
        src = os.path.join(TMP, "pp.pptx")
        prs = Presentation(); prs.slides.add_slide(prs.slide_layouts[1])
        prs.save(src)
        assert DocumentConverter().convert(src, os.path.join(TMP, "pp.txt"))


class TestImage:
    def test_png_jpg(self):
        from core.image_converter import ImageConverter
        src = os.path.join(TMP, "i.png"); _make_png(src)
        out = os.path.join(TMP, "i.jpg")
        assert ImageConverter().convert(src, out)
        assert os.path.isfile(out)

    def test_compress(self):
        from core.tools import image_compress
        src = os.path.join(TMP, "c.png"); _make_png(src, (640, 480))
        out = os.path.join(TMP, "c.jpg")
        assert image_compress(src, out, quality=60)

    def test_watermark(self):
        from core.watermark_tool import process_watermark
        src = os.path.join(TMP, "w.png"); _make_png(src)
        assert process_watermark(src, os.path.join(TMP, "w_out.png"), "text", "水印")

    def test_image_merge(self):
        from core.image_album import merge_vertical
        a = os.path.join(TMP, "m1.png"); _make_png(a)
        b = os.path.join(TMP, "m2.png"); _make_png(b, color=(50, 150, 200))
        assert merge_vertical([a, b], os.path.join(TMP, "m.png"))

    def test_phantom(self):
        from core.phantom_tank import make_phantom
        w = os.path.join(TMP, "ph_w.png")
        _make_png(w, size=(100, 100), color=(255, 255, 255))
        b = os.path.join(TMP, "ph_b.png")
        _make_png(b, size=(100, 100), color=(0, 0, 0))
        assert make_phantom(w, b, os.path.join(TMP, "ph.png"))

    def test_id_photo(self):
        from core.id_photo import change_background
        src = os.path.join(TMP, "id.png")
        _make_png(src, size=(300, 400), color=(240, 200, 180))
        out = os.path.join(TMP, "id_blue.jpg")
        assert change_background(src, out, "蓝底", use_ai=False)


class TestAudio:
    def test_wav_mp3(self):
        import struct
        import wave
        from core.audio_converter import AudioConverter
        src = os.path.join(TMP, "a.wav")
        with wave.open(src, "wb") as w:
            w.setnchannels(1); w.setsampwidth(2); w.setframerate(8000)
            w.writeframes(b"".join(struct.pack("<h", 10000) for _ in range(8000)))
        out = os.path.join(TMP, "a.mp3")
        assert AudioConverter().convert(src, out, codec="libmp3lame")
        assert os.path.isfile(out)


class TestVideo:
    def test_mp4_avi(self):
        from core.video_converter import VideoConverter
        src = os.path.join(TMP, "v.mp4"); _make_video(src)
        assert VideoConverter().convert(src, os.path.join(TMP, "v.avi"),
                                        ".avi", "libx264", "medium")

    def test_gif_encode(self):
        """GIF 编码能力（面板 runner 已单独验证，此处直接测 ffmpeg 转 gif）。"""
        src = os.path.join(TMP, "g.mp4"); _make_video(src, duration=1, size="100x80", rate=8)
        out = os.path.join(TMP, "g.gif")
        ff = _ff()
        r = subprocess.run([ff, "-y", "-i", src, "-vf", "fps=10", "-loop", "0", out],
                           capture_output=True, timeout=60)
        assert r.returncode == 0 and os.path.isfile(out) and os.path.getsize(out) > 0

    def test_thumbnail_sheet(self):
        from core.thumbnail_sheet import generate_thumbnail_sheet
        src = os.path.join(TMP, "t.mp4"); _make_video(src, duration=2)
        out = os.path.join(TMP, "sheet.jpg")
        assert generate_thumbnail_sheet(src, out, cols=2, rows=2)


class TestRecognition:
    def test_ocr(self):
        from core.ocr_tool import ocr_image
        src = os.path.join(TMP, "ocr.png")
        _make_png(src, (320, 80), text="Test OCR 2026")
        text = ocr_image(src)
        assert text and "OCR" in text

    def test_table_csv(self):
        from core.table_recognizer import table_to_csv
        from PIL import Image, ImageDraw
        src = os.path.join(TMP, "tbl.png")
        img = Image.new("RGB", (500, 160), "white")
        d = ImageDraw.Draw(img)
        for x in (20, 260, 480):
            d.line([(x, 20), (x, 140)], fill="black", width=2)
        for y in (20, 80, 140):
            d.line([(20, y), (480, y)], fill="black", width=2)
        d.text((50, 45), "Name Value", fill="black")
        img.save(src)
        out = os.path.join(TMP, "tbl.csv")
        assert table_to_csv(src, out)
        assert os.path.isfile(out)

    def test_mediainfo(self):
        from core.mediainfo import get_mediainfo
        src = os.path.join(TMP, "mi.png"); _make_png(src)
        assert get_mediainfo(src) is not None


class TestUtils:
    def test_hash(self):
        from core.hash_tool import compute_hash
        src = os.path.join(TMP, "h.txt")
        open(src, "w", encoding="utf-8").write("abc")
        h = compute_hash(src, "SHA256")
        assert h and len(h) == 64

    def test_qrcode(self):
        import qrcode
        out = os.path.join(TMP, "qr.png")
        qr = qrcode.QRCode(version=None, error_correction=qrcode.constants.ERROR_CORRECT_H)
        qr.add_data("https://example.com"); qr.make(fit=True)
        qr.make_image().save(out)
        assert os.path.getsize(out) > 100

    def test_batch_rename(self):
        from core.tools import batch_rename
        files = []
        for i in range(3):
            p = os.path.join(TMP, f"old_{i}.txt")
            open(p, "w", encoding="utf-8").write(str(i))
            files.append(p)
        out_dir = os.path.join(TMP, "renamed")
        os.makedirs(out_dir, exist_ok=True)
        assert batch_rename(files, "new_{n}", 1, output_dir=out_dir)
        assert len(os.listdir(out_dir)) == 3


class TestSecurity:
    def test_aes_roundtrip(self):
        from core.file_security import encrypt_file, decrypt_file
        src = os.path.join(TMP, "sec.txt")
        open(src, "w", encoding="utf-8").write("机密 TOP SECRET")
        enc = os.path.join(TMP, "sec.enc")
        dec = os.path.join(TMP, "sec_dec.txt")
        assert encrypt_file(src, enc, "pwd")
        assert decrypt_file(enc, dec, "pwd")
        assert open(dec, encoding="utf-8").read() == "机密 TOP SECRET"

    def test_shred(self):
        from core.file_security import shred_file
        src = os.path.join(TMP, "shred.txt")
        open(src, "w", encoding="utf-8").write("x")
        assert shred_file(src)
        assert not os.path.exists(src)


class TestTaskFlow:
    def test_task_executes(self):
        """真实任务链路：add_task -> 执行 -> SUCCESS。"""
        os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
        from PySide6.QtWidgets import QApplication
        app = QApplication.instance() or QApplication([])

        from gui_qt import task_manager as tm
        from gui_qt.components.theme_manager import ThemeManager
        from gui_qt.services import QtServices
        from gui_qt.task_manager import TaskManager
        services = QtServices()
        services.task_manager = TaskManager(services, None)
        services.theme_mgr = ThemeManager(services)
        mgr = services.task_manager

        src = os.path.join(TMP, "flow.txt")
        open(src, "w", encoding="utf-8").write("flow test")
        out = os.path.join(TMP, "flow.docx")

        def runner(task, prog):
            from core.doc_converter import DocumentConverter
            return DocumentConverter().convert(task.file_path, task.output_path, prog)

        tid = mgr.add_task(name="t", task_type="doc", file_path=src,
                           output_path=out, params={}, runner=runner)
        assert tid is not None

        import time
        deadline = time.time() + 30
        while time.time() < deadline:
            app.processEvents()
            t = mgr.get_task(tid)
            if t and t.state in (tm.SUCCESS, tm.FAILED, tm.CANCELLED):
                break
            time.sleep(0.05)
        t = mgr.get_task(tid)
        assert t.state == tm.SUCCESS, f"state={t.state} error={t.error}"
        assert os.path.isfile(out)


    # ── 统一服务不含转换 API（已去除）──────────
    def test_lan_service_no_api(self):
        """统一局域网服务不再暴露 /api 转换路由。"""
        from core.lan_service import LanService
        srv = LanService(host="127.0.0.1", port=18779)
        import tempfile
        d = tempfile.mkdtemp()
        srv.set_recv(d)
        assert srv.start()
        try:
            import urllib.request
            import urllib.error
            base = f"http://127.0.0.1:{srv.port}"
            # /api/health 应 404（API 已去除）
            try:
                with urllib.request.urlopen(base + "/api/health", timeout=5) as r:
                    assert r.status != 200, "API 路由仍存在"
            except urllib.error.HTTPError as e:
                assert e.code in (404, 405), f"意外状态 {e.code}"
            # 首页不应含「转换 API」卡片
            with urllib.request.urlopen(base + "/", timeout=5) as r:
                assert "转换 API" not in r.read().decode()
        finally:
            srv.stop()


class TestLanService:
    """合并后的统一局域网服务（分享 + 接收 + API 单端口）。"""

    def _mk(self):
        from core.lan_service import LanService
        return LanService(host="127.0.0.1", port=18941)

    def test_share_download(self):
        srv = self._mk()
        share_f = os.path.join(TMP, "ls_share.txt")
        open(share_f, "w", encoding="utf-8").write("LAN SHARE test")
        srv.start_share([share_f])
        assert srv.start()
        try:
            import urllib.request
            base = f"http://127.0.0.1:{srv.port}"
            with urllib.request.urlopen(base + "/", timeout=5) as r:
                assert "文件分享" in r.read().decode()
            with urllib.request.urlopen(base + "/share/", timeout=5) as r:
                assert "ls_share.txt" in r.read().decode()
            dl = []
            srv.on_downloaded = lambda n, s, t, ip: dl.append(n)
            with urllib.request.urlopen(base + "/share/ls_share.txt", timeout=5) as r:
                assert r.read().decode() == "LAN SHARE test"
            assert dl == ["ls_share.txt"]
        finally:
            srv.stop()

    def test_recv_upload(self):
        srv = self._mk()
        recv_dir = os.path.join(TMP, "ls_recv")
        srv.set_recv(recv_dir)
        assert srv.start()
        try:
            import urllib.request
            base = f"http://127.0.0.1:{srv.port}"
            boundary = "----fmtst"
            body = (f'--{boundary}\r\nContent-Disposition: form-data; name="file"; '
                    f'filename="up.txt"\r\nContent-Type: text/plain\r\n\r\n'
                    f"UP CONTENT\r\n--{boundary}--\r\n").encode()
            req = urllib.request.Request(
                base + "/recv/upload", data=body,
                headers={"Content-Type": f"multipart/form-data; boundary={boundary}"})
            got = []
            srv.on_received = lambda n, s, t, ip, rn=None: got.append(n)
            with urllib.request.urlopen(req, timeout=10) as r:
                assert r.status == 200
            assert got == ["up.txt"]
            assert os.path.isfile(os.path.join(recv_dir, "up.txt"))
        finally:
            srv.stop()
