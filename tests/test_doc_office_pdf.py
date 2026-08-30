"""doc_office_pdf 多级降级引擎单元测试。

覆盖:
- 纯 Python 渲染(DOCX/PPTX → PDF,含中文与表格内容验证)
- 降级链顺序(mock 各引擎,验证优先级与兜底)
- 旧格式 .doc/.ppt 的明确错误提示
运行:venv/Scripts/python -m pytest tests/test_doc_office_pdf.py -q
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pytest


@pytest.fixture(scope="module")
def docx_path(tmp_path_factory):
    """构造含标题/段落/表格的测试 docx。"""
    from docx import Document
    d = Document()
    d.add_heading("测试标题", level=1)
    d.add_paragraph("中文字体测试段落：格式大师降级引擎。")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "姓名"; t.cell(0, 1).text = "分数"
    t.cell(1, 0).text = "张三"; t.cell(1, 1).text = "95"
    p = tmp_path_factory.mktemp("fm") / "test.docx"
    d.save(str(p))
    return str(p)


@pytest.fixture(scope="module")
def pptx_path(tmp_path_factory):
    """构造含文本框的测试 pptx。"""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tb.text_frame.text = "PPT 降级引擎测试"
    p = tmp_path_factory.mktemp("fm") / "test.pptx"
    prs.save(str(p))
    return str(p)


def _pdf_text(pdf_path):
    """提取 PDF 全部文本用于断言。"""
    import pymupdf
    d = pymupdf.open(pdf_path)
    try:
        return "".join(p.get_text() for p in d)
    finally:
        d.close()


class TestPureRender:
    def test_docx_to_pdf_content(self, docx_path, tmp_path):
        from core.doc_office_pdf import render_docx_to_pdf
        out = str(tmp_path / "out.pdf")
        assert render_docx_to_pdf(docx_path, out) is True
        text = _pdf_text(out)
        assert "测试标题" in text
        assert "格式大师" in text
        assert "张三" in text  # 表格内容
        assert "95" in text

    def test_pptx_to_pdf_content(self, pptx_path, tmp_path):
        from core.doc_office_pdf import render_pptx_to_pdf
        out = str(tmp_path / "out.pdf")
        assert render_pptx_to_pdf(pptx_path, out) is True
        text = _pdf_text(out)
        assert "PPT 降级引擎测试" in text

    def test_xlsx_converter_preserves_chinese_text(self, tmp_path, monkeypatch):
        """文档面板的 XLSX→PDF 入口也必须走 CJK 引擎。"""
        import openpyxl
        import core.doc_office_pdf as office_pdf
        from core.doc_converter import DocumentConverter

        source = tmp_path / "积分清单.xlsx"
        output = tmp_path / "积分清单.pdf"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.append(["手机号", "惠通积分", "微盟积分", "差异"])
        sheet.append(["13800138000", 210, 200, 10])
        workbook.save(source)
        monkeypatch.setattr(office_pdf, "_soffice_available", lambda: None)

        assert DocumentConverter().convert(str(source), str(output)) is True
        text = _pdf_text(str(output))
        assert "惠通积分" in text
        assert "微盟积分" in text
        assert "差异" in text


class TestFallbackChain:
    def test_docx_all_engines_down_uses_builtin(self, docx_path, tmp_path, monkeypatch):
        """COM 与 LibreOffice 全部不可用时，DOCX 降级到内置引擎。"""
        import core.doc_office_pdf as m
        monkeypatch.setattr(m, "_com_dispatch", lambda progids: (None, None))
        monkeypatch.setattr(m, "_soffice_available", lambda: None)
        out = str(tmp_path / "fallback.pdf")
        ok, engine = m.docx_to_pdf(docx_path, out)
        assert ok is True
        assert engine == "内置引擎"
        assert os.path.isfile(out) and os.path.getsize(out) > 0

    def test_pptx_all_engines_down_uses_builtin(self, pptx_path, tmp_path, monkeypatch):
        import core.doc_office_pdf as m
        monkeypatch.setattr(m, "_com_dispatch", lambda progids: (None, None))
        monkeypatch.setattr(m, "_soffice_available", lambda: None)
        out = str(tmp_path / "fallback.pdf")
        ok, engine = m.ppt_to_pdf(pptx_path, out)
        assert ok is True
        assert engine == "内置引擎"

    def test_doc_legacy_needs_office(self, tmp_path, monkeypatch):
        """旧 .doc 二进制格式：无办公软件时给出明确错误。"""
        import core.doc_office_pdf as m
        monkeypatch.setattr(m, "_com_dispatch", lambda progids: (None, None))
        monkeypatch.setattr(m, "_soffice_available", lambda: None)
        fake = str(tmp_path / "legacy.doc")
        open(fake, "wb").write(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 100)
        with pytest.raises(RuntimeError) as ei:
            m.docx_to_pdf(fake, str(tmp_path / "out.pdf"))
        assert "需本机安装" in str(ei.value)

    def test_ppt_legacy_needs_office(self, tmp_path, monkeypatch):
        import core.doc_office_pdf as m
        monkeypatch.setattr(m, "_com_dispatch", lambda progids: (None, None))
        monkeypatch.setattr(m, "_soffice_available", lambda: None)
        fake = str(tmp_path / "legacy.ppt")
        open(fake, "wb").write(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 100)
        with pytest.raises(RuntimeError) as ei:
            m.ppt_to_pdf(fake, str(tmp_path / "out.pdf"))
        assert "需本机安装" in str(ei.value)


class TestLibreOffice:
    def test_find_soffice_candidates(self, monkeypatch):
        """soffice 探测逻辑：常见路径 + 环境变量 + PATH。"""
        import core.doc_office_pdf as m
        # 路径不存在 → 返回 None（CI 环境）
        monkeypatch.setenv("LIBREOFFICE_PATH", "")
        r = m._find_soffice()
        assert r is None or os.path.isfile(r)

    def test_register_cjk_font_fallback(self):
        """字体注册必须有结果（项目字体或系统字体或 Helvetica）。"""
        from core.doc_office_pdf import _register_cjk_font
        name, is_cjk = _register_cjk_font()
        assert name  # 非空即注册成功


class TestChineseTextEncoding:
    """本地常见的 GBK 文本与 CSV 不得在转换后出现乱码。"""

    def test_gbk_txt_to_docx(self, tmp_path):
        from docx import Document
        from core.doc_converter import DocumentConverter

        source = tmp_path / "中文备注.txt"
        output = tmp_path / "中文备注.docx"
        source.write_bytes("积分对账调整清单\n客户：北海".encode("gbk"))

        assert DocumentConverter().convert(str(source), str(output)) is True
        text = "\n".join(p.text for p in Document(output).paragraphs)
        assert "积分对账调整清单" in text
        assert "客户：北海" in text

    def test_gbk_csv_to_xlsx(self, tmp_path):
        import openpyxl
        from core.doc_converter import DocumentConverter

        source = tmp_path / "对账.csv"
        output = tmp_path / "对账.xlsx"
        source.write_bytes("姓名,积分\n北海,210\n".encode("gbk"))

        assert DocumentConverter().convert(str(source), str(output)) is True
        sheet = openpyxl.load_workbook(output).active
        assert sheet["A1"].value == "姓名"
        assert sheet["A2"].value == "北海"
