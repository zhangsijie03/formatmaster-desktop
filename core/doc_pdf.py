"""PDF 格式转换"""
import os


class DocPdfMixin:
    """PDF 转换方法"""

    # ========== PDF 转换 ==========

    def _pdf_to_docx(self, inp, out, cb):
        if cb: cb(20, "解析PDF...")
        # pdf2docx 内部 `import fitz`，会触发 PyMuPDF 弃用提示打印（stderr 噪音）。
        # 预先把 fitz 别名绑定到 pymupdf：绕过 fitz/__init__.py 的 message_warning，
        # 同时保证 pdf2docx 拿到的就是完整 pymupdf 模块（功能完全等价）。
        import sys
        import pymupdf
        sys.modules.setdefault("fitz", pymupdf)
        from pdf2docx import Converter
        cv = Converter(inp)
        if cb: cb(50, "生成Word...")
        cv.convert(out)
        cv.close()
        return True

    def _pdf_to_txt(self, inp, out, cb):
        import pymupdf
        doc = pymupdf.open(inp)
        text = []
        total = len(doc)
        for i, page in enumerate(doc):
            if self._cancel: return False
            text.append(page.get_text())
            if cb:
                cb(20 + int(i * 70 / max(total, 1)), f"读取第{i+1}/{total}页...")
        doc.close()
        self._write_text(out, '\n'.join(text))
        return True

    def _pdf_to_image(self, inp, out, cb):
        import pymupdf
        doc = pymupdf.open(inp)
        total = len(doc)
        out_dir = os.path.dirname(out)
        base = os.path.splitext(os.path.basename(out))[0]
        ext = os.path.splitext(out)[1]
        for i, page in enumerate(doc):
            if self._cancel: return False
            pix = page.get_pixmap(dpi=200)
            if total == 1:
                pix.save(out)
            else:
                pix.save(os.path.join(out_dir, f"{base}_{i+1}{ext}"))
            if cb:
                cb(20 + int(i * 70 / max(total, 1)), f"渲染第{i+1}/{total}页...")
        doc.close()
        return True

    def _pdf_to_html(self, inp, out, cb):
        import pymupdf
        doc = pymupdf.open(inp)
        if cb: cb(30, "转换中...")
        parts = ["<html><head><meta charset='utf-8'></head><body>"]
        for page in doc:
            parts.append(page.get_text("html"))
        parts.append("</body></html>")
        doc.close()
        self._write_text(out, '\n'.join(parts))
        return True

    def _pdf_to_pptx(self, inp, out, cb):
        import pymupdf
        from pptx import Presentation
        from pptx.util import Inches
        doc = pymupdf.open(inp)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        total = len(doc)
        for i, page in enumerate(doc):
            if self._cancel: return False
            text = page.get_text().strip()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            from pptx.util import Pt
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            if text:
                for line in text.split('\n')[:50]:
                    p = tf.add_paragraph()
                    p.text = line[:200]
                    p.font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = f"(第{i+1}页 - 无文字内容)"
                p.font.size = Pt(18)
            if cb:
                cb(20 + int(i * 70 / max(total, 1)), f"第{i+1}/{total}页...")
        doc.close()
        prs.save(out)
        return True

    def _pdf_to_xlsx(self, inp, out, cb):
        import pymupdf
        import openpyxl
        doc = pymupdf.open(inp)
        wb = openpyxl.Workbook()
        ws = wb.active
        total = len(doc)
        for i, page in enumerate(doc):
            if self._cancel: return False
            text = page.get_text().strip().split('\n')
            for j, line in enumerate(text):
                ws.cell(row=i+1, column=j+1, value=line[:32767])
            if cb:
                cb(20 + int(i * 70 / max(total, 1)), f"第{i+1}/{total}页...")
        doc.close()
        wb.save(out)
        return True
