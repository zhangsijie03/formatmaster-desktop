"""Word/WPS 格式转换"""
import os


class DocWordMixin:
    """Word/WPS 转换方法"""

    # ========== Word 转换 ==========

    def _docx_to_pdf(self, inp, out, cb):
        """DOCX/DOC/WPS → PDF（多级降级：Word COM → WPS COM → LibreOffice → 内置引擎）。

        优先使用本机 Office 保真排版；无办公软件时自动降级内置引擎，
        不再因缺少 Office 直接失败（.doc 旧格式除外，需办公软件）。
        """
        if cb: cb(30, "准备导出PDF...")
        from core.doc_office_pdf import docx_to_pdf
        try:
            ok, engine = docx_to_pdf(inp, out, cb)
            if cb:
                cb(100, f"转换完成（{engine}）")
            return ok
        except Exception as e:
            if cb:
                cb(-1, f"错误: {str(e)[:200]}")
            return False

    def _docx_to_txt(self, inp, out, cb):
        if cb: cb(30, "读取文档...")
        import docx
        doc = docx.Document(inp)
        text = '\n'.join(p.text for p in doc.paragraphs)
        self._write_text(out, text)
        return True

    def _docx_to_html(self, inp, out, cb):
        if cb: cb(30, "转换中...")
        import docx
        doc = docx.Document(inp)
        body_parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                body_parts.append('<br>')
                continue
            tag = 'p'
            if para.style and para.style.name:
                sn = para.style.name.lower()
                if 'heading 1' in sn or '标题 1' in sn: tag = 'h1'
                elif 'heading 2' in sn or '标题 2' in sn: tag = 'h2'
                elif 'heading 3' in sn or '标题 3' in sn: tag = 'h3'
            body_parts.append(f'<{tag}>{self._safe_html(text)}</{tag}>')
        if doc.tables:
            for table in doc.tables:
                body_parts.append('<table border="1" cellpadding="4" style="border-collapse:collapse">')
                for row in table.rows:
                    body_parts.append('<tr>')
                    for cell in row.cells:
                        body_parts.append(f'<td>{self._safe_html(cell.text.strip())}</td>')
                    body_parts.append('</tr>')
                body_parts.append('</table><br>')
        self._write_text(out, self._build_html_page('\n'.join(body_parts)))
        return True

    def _docx_to_image(self, inp, out, cb):
        if cb: cb(20, "读取文档...")
        import docx
        from PIL import Image, ImageDraw, ImageFont
        doc = docx.Document(inp)
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        if not lines:
            lines = ["(空白文档)"]
        font_size = 20
        try:
            font = ImageFont.truetype("msyh.ttc", font_size)
        except Exception:
            font = ImageFont.load_default()
        line_h = font_size + 10
        img_w = 1000
        img_h = max(200, len(lines) * line_h + 100)
        img = Image.new('RGB', (img_w, img_h), 'white')
        draw = ImageDraw.Draw(img)
        y = 20
        for line in lines:
            draw.text((20, y), line[:120], fill='black', font=font)
            y += line_h
        img.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _docx_to_pptx(self, inp, out, cb):
        if cb: cb(20, "读取文档...")
        import docx
        from pptx import Presentation
        from pptx.util import Inches, Pt
        doc = docx.Document(inp)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        chunk_size = 15
        for i in range(0, max(len(lines), 1), chunk_size):
            if self._cancel: return False
            chunk = lines[i:i+chunk_size]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            if chunk:
                for j, line in enumerate(chunk):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line[:200]
                    p.font.size = Pt(18)
            if cb:
                cb(20 + int(i * 70 / max(len(lines), 1)), f"段落{i+1}...")
        prs.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _docx_to_md(self, inp, out, cb):
        if cb: cb(30, "读取文档...")
        import docx
        doc = docx.Document(inp)
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                lines.append('')
                continue
            if para.style and para.style.name:
                sn = para.style.name.lower()
                if 'heading 1' in sn or '标题 1' in sn:
                    lines.append(f'# {text}')
                    continue
                elif 'heading 2' in sn or '标题 2' in sn:
                    lines.append(f'## {text}')
                    continue
                elif 'heading 3' in sn or '标题 3' in sn:
                    lines.append(f'### {text}')
                    continue
            lines.append(text)
        self._write_text(out, '\n\n'.join(lines))
        if cb: cb(100, "转换完成")
        return True

    def _doc_copy(self, inp, out, cb):
        if cb: cb(50, "复制转换中...")
        import shutil
        shutil.copy2(inp, out)
        return True
