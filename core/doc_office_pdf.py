"""doc_office_pdf — 文档→PDF 多级降级引擎（不依赖本机 Office 也能出 PDF）。

降级链（按优先级自动尝试，任一成功即返回）：
1. Microsoft Office COM（Word / PowerPoint，质量最高，排版保真）
2. WPS Office COM（国内常见，接口兼容 Word/PowerPoint）
3. LibreOffice headless（soffice --headless --convert-to pdf）
4. 纯 Python 渲染（reportlab，读取 docx/pptx 内容排版；不依赖任何办公软件）

设计要点：
- 每级失败静默降级，最后一级失败才抛错（错误信息含各引擎失败原因）；
- COM 操作带超时保护：每次调用包 try/except，超时/异常都视为引擎不可用；
- 纯 Python 渲染用系统中文字体（微软雅黑/宋体/黑体）降级注册；
- 返回 (ok, message)；message 用于进度回调展示当前引擎，方便用户理解。
"""
import os
import shutil
import subprocess
import sys

# ── 引擎探测（进程级缓存，避免重复探测）────────────────────
_detected = {}


def _sys32():
    """System32 路径（win32com 需要加载的 COM 环境无需额外处理）。"""
    return os.environ.get("SystemRoot", "C:/Windows") + "/System32"


def _find_soffice():
    """查找 LibreOffice 可执行文件（常见安装路径）。"""
    candidates = [os.environ.get("LIBREOFFICE_PATH", "")]
    if os.name == "nt":
        candidates.extend([
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"C:\Program Files\LibreOffice\program\soffice.com",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.com",
        ])
    elif sys.platform == "darwin":
        candidates.extend([
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            os.path.expanduser(
                "~/Applications/LibreOffice.app/Contents/MacOS/soffice"),
            "/opt/homebrew/bin/soffice",
            "/usr/local/bin/soffice",
        ])
    for c in candidates:
        if c and os.path.isfile(c):
            return c
    # PATH 中的 soffice
    found = shutil.which("soffice") or shutil.which("soffice.exe")
    return found or None


def _soffice_available():
    if "soffice" not in _detected:
        _detected["soffice"] = _find_soffice()
    return _detected["soffice"]


# ── COM 分发（Word / PowerPoint / WPS）────────────────────

def _com_dispatch(progids):
    """尝试按顺序分发 COM 应用；成功返回 (app, progid_used)，全失败返回 (None, None)。

    注意：调用方负责在成功后 _com_cleanup(app)（内含 Quit + CoUninitialize）；
    全部失败时本函数自行 CoUninitialize，避免线程 COM 状态残留。
    """
    if os.name != "nt":
        return None, None
    import pythoncom
    pythoncom.CoInitialize()
    import win32com.client
    for progid in progids:
        try:
            app = win32com.client.DispatchEx(progid)
            return app, progid
        except Exception:
            continue
    try:
        pythoncom.CoUninitialize()
    except Exception:
        pass
    return None, None


def _com_cleanup(app):
    """安全关闭 COM 应用实例。"""
    if app is None:
        return
    try:
        app.Quit()
    except Exception:
        pass
    try:
        import pythoncom
        pythoncom.CoUninitialize()
    except Exception:
        pass


def _word_com(app, inp, out):
    """用 Word/WPS COM 把 docx 转 PDF；成功返回 True。"""
    doc = None
    try:
        app.Visible = False
        try:
            app.DisplayAlerts = 0
        except Exception:
            pass
        doc = app.Documents.Open(os.path.abspath(inp), ReadOnly=True)
        doc.SaveAs2(os.path.abspath(out), FileFormat=17)  # 17 = wdFormatPDF
        return True
    finally:
        if doc is not None:
            try:
                doc.Close(0)
            except Exception:
                pass


def _ppt_com(app, inp, out):
    """用 PowerPoint/WPS COM 把 ppt/pptx 转 PDF；成功返回 True。"""
    pres = None
    try:
        app.Visible = False
        pres = app.Presentations.Open(os.path.abspath(inp), WithWindow=False)
        pres.SaveAs(os.path.abspath(out), 32)  # 32 = ppSaveAsPDF
        return True
    finally:
        if pres is not None:
            try:
                pres.Close()
            except Exception:
                pass


# ── LibreOffice headless ────────────────────────────────

def _libreoffice_convert(inp, out):
    """用 LibreOffice headless 转换；成功返回 True。"""
    soffice = _soffice_available()
    if not soffice:
        return False
    out_dir = os.path.dirname(os.path.abspath(out)) or "."
    cmd = [
        soffice, "--headless", "--convert-to", "pdf",
        "--outdir", out_dir, os.path.abspath(inp),
    ]
    # CREATE_NO_WINDOW：避免转换时弹出黑色控制台
    creationflags = 0x08000000 if os.name == "nt" else 0
    try:
        r = subprocess.run(cmd, capture_output=True, timeout=120,
                           creationflags=creationflags,
                           encoding="utf-8", errors="ignore")
        # LibreOffice 输出文件名 = 输入基名.pdf（在 outdir）
        base = os.path.splitext(os.path.basename(inp))[0]
        expect = os.path.join(out_dir, base + ".pdf")
        if os.path.isfile(expect):
            # 若目标名不同（面板可能带 _converted 后缀），移动过去
            if os.path.abspath(expect) != os.path.abspath(out):
                shutil.move(expect, out)
            return True
        # 直接输出到目标路径的情况
        return r.returncode == 0 and os.path.isfile(out)
    except Exception:
        return False


# ── 纯 Python 渲染（最终兜底，无需任何办公软件）────────────

def _register_cjk_font():
    """注册中文字体，返回 (font_name, is_cjk)。用系统默认字体。"""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfbase.ttfonts import TTFont
    candidates = [
        "C:/Windows/Fonts/msyh.ttc",   # 微软雅黑（系统默认）
        "C:/Windows/Fonts/simsun.ttc",  # 宋体
        "C:/Windows/Fonts/simhei.ttf",  # 黑体
    ]
    if sys.platform == "darwin":
        candidates = [
            "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
            "/System/Library/Fonts/Supplemental/Arial.ttf",
            "/System/Library/Fonts/Supplemental/Songti.ttc",
            "/System/Library/Fonts/Supplemental/PingFang.ttc",
            "/System/Library/Fonts/PingFang.ttc",
        ] + candidates
    for fp in candidates:
        try:
            if os.path.isfile(fp):
                pdfmetrics.registerFont(TTFont("CJK", fp))
                return "CJK", True
        except Exception:
            continue
    # 精简 Linux/容器通常没有 CJK 系统字体。ReportLab 自带的 CID 字体
    # 可保证中文可见且保留 PDF 文本层，避免降级成不可搜索的方框字符。
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        return "STSong-Light", True
    except Exception:
        return "Helvetica", False


def render_docx_to_pdf(inp, out, progress_cb=None):
    """纯 Python：docx → PDF（reportlab，段落/标题/表格/图片）。

    不依赖 Office/WPS/LibreOffice。排版为近似效果（保文本、标题层级、
    简单表格），复杂版式（页眉页脚/分栏/嵌入字体）会简化。
    """
    from io import BytesIO
    from docx import Document
    from docx.document import Document as _Doc
    from docx.oxml.ns import qn
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Table, TableStyle, Image as RLImage)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import mm

    font_name, _ = _register_cjk_font()
    if progress_cb:
        progress_cb(20, "读取文档内容...")

    doc = Document(inp)
    pdf = SimpleDocTemplate(out, pagesize=A4,
                            leftMargin=20 * mm, rightMargin=20 * mm,
                            topMargin=18 * mm, bottomMargin=18 * mm)
    styles = getSampleStyleSheet()
    body = ParagraphStyle("Body", parent=styles["Normal"],
                          fontName=font_name, fontSize=10.5, leading=16)
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontName=font_name)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontName=font_name)
    h3 = ParagraphStyle("H3", parent=styles["Heading3"], fontName=font_name)
    story = []

    def _safe(text):
        return (text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    def _walk_element(el):
        """递归处理文档元素（段落/标题/表格/图片）。"""
        if isinstance(el, _Para):
            text = el.text.strip()
            # 段落内嵌图片（先于文本输出，保持大致顺序）
            try:
                for run in el.runs:
                    for blip in run._element.findall(".//" + qn("a:blip")):
                        rid = blip.get(qn("r:embed"))
                        if not rid:
                            continue
                        try:
                            part = doc.part.related_parts[rid]
                            blob = part.blob
                            img = RLImage(BytesIO(blob))
                            # 限制宽度到内容区（A4 - 40mm 边距 ≈ 170mm）
                            avail_w = 170 * mm
                            iw, ih = img.imageWidth, img.imageHeight
                            if iw > avail_w:
                                scale = avail_w / float(iw)
                                img.drawWidth = avail_w
                                img.drawHeight = ih * scale
                            story.append(img)
                            story.append(Spacer(1, 6))
                        except Exception:
                            continue
            except Exception:
                pass
            if not text:
                story.append(Spacer(1, 8))
                return
            style_name = (el.style.name or "").lower() if el.style else ""
            if "heading 1" in style_name or "标题 1" in style_name:
                story.append(Paragraph(_safe(text), h1))
            elif "heading 2" in style_name or "标题 2" in style_name:
                story.append(Paragraph(_safe(text), h2))
            elif "heading 3" in style_name or "标题 3" in style_name:
                story.append(Paragraph(_safe(text), h3))
            else:
                story.append(Paragraph(_safe(text), body))
            story.append(Spacer(1, 4))
        elif isinstance(el, _Tbl):
            rows = []
            for row in el.rows:
                rows.append([_safe(c.text) for c in row.cells])
            if rows:
                t = Table(rows, repeatRows=1)
                t.setStyle(TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.93, 0.93, 0.93)),
                    ("FONTNAME", (0, 0), (-1, -1), font_name),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ]))
                story.append(t)
                story.append(Spacer(1, 10))

    from docx.table import Table as _Tbl
    from docx.text.paragraph import Paragraph as _Para
    total = len(doc.paragraphs) + len(doc.tables)
    processed = 0
    # 遍历 body 元素保持文档顺序
    for el in doc.element.body:
        tag = el.tag
        try:
            if tag == qn("w:p"):
                para = _Para(el, doc)
                _walk_element(para)
            elif tag == qn("w:tbl"):
                tbl = _Tbl(el, doc)
                _walk_element(tbl)
        except Exception:
            pass
        processed += 1
        if progress_cb and processed % 5 == 0:
            pct = 30 + int(processed * 55 / max(total, 1))
            progress_cb(min(pct, 85), "排版中...")

    if progress_cb:
        progress_cb(88, "生成PDF...")
    pdf.build(story)
    return os.path.isfile(out)


def render_pptx_to_pdf(inp, out, progress_cb=None):
    """纯 Python：pptx → PDF（reportlab，每页一张幻灯片，文本排版）。

    不依赖 Office。老版 .ppt 二进制格式 python-pptx 无法读取，返回 False。
    """
    from pptx import Presentation
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm

    font_name, _ = _register_cjk_font()
    if progress_cb:
        progress_cb(20, "读取演示文稿...")

    prs = Presentation(inp)
    pdf = SimpleDocTemplate(out, pagesize=landscape(A4),
                            leftMargin=12 * mm, rightMargin=12 * mm,
                            topMargin=12 * mm, bottomMargin=12 * mm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("SlideTitle", parent=styles["Heading1"],
                                 fontName=font_name, fontSize=16, spaceAfter=8)
    body_style = ParagraphStyle("SlideBody", parent=styles["Normal"],
                                fontName=font_name, fontSize=11, leading=15)
    story = []

    def _safe(text):
        return (text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    for i, slide in enumerate(prs.slides):
        story.append(Paragraph(f"幻灯片 {i + 1}", title_style))
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                texts.append(shape.text.strip())
        for t in texts[:60]:
            story.append(Paragraph(_safe(t), body_style))
        story.append(Spacer(1, 18))
        if progress_cb:
            pct = 30 + int(i * 60 / max(len(prs.slides), 1))
            progress_cb(min(pct, 88), f"幻灯片{i+1}...")
    if progress_cb:
        progress_cb(90, "生成PDF...")
    pdf.build(story)
    return os.path.isfile(out)


# ── 高层统一入口 ─────────────────────────────────────────

def docx_to_pdf(inp, out, progress_cb=None):
    """DOCX/DOC/WPS → PDF 多级降级。返回 (ok, message)。"""
    msgs = []
    # 1) Microsoft Word COM
    try:
        app, progid = _com_dispatch(["Word.Application"])
        if app is not None:
            if progress_cb:
                progress_cb(40, "通过 Word 导出 PDF...")
            ok = _word_com(app, inp, out)
            _com_cleanup(app)
            if ok and os.path.isfile(out):
                return True, "Word"
            msgs.append("Word COM 导出失败")
    except Exception as e:
        msgs.append(f"Word 不可用({type(e).__name__})")
    # 2) WPS 文字 COM
    try:
        app, progid = _com_dispatch(["Kwps.Application", "wps.Application"])
        if app is not None:
            if progress_cb:
                progress_cb(40, "通过 WPS 导出 PDF...")
            ok = _word_com(app, inp, out)
            _com_cleanup(app)
            if ok and os.path.isfile(out):
                return True, "WPS"
            msgs.append("WPS COM 导出失败")
    except Exception as e:
        msgs.append(f"WPS 不可用({type(e).__name__})")
    # 3) LibreOffice headless
    if _soffice_available():
        if progress_cb:
            progress_cb(40, "通过 LibreOffice 导出 PDF...")
        if _libreoffice_convert(inp, out):
            return True, "LibreOffice"
        msgs.append("LibreOffice 转换失败")
    else:
        msgs.append("未安装 LibreOffice")
    # 4) 纯 Python（仅 docx 可读；老版 .doc 二进制无法用 python-docx 读取）
    ext = os.path.splitext(inp)[1].lower()
    if ext in (".docx", ".wps"):
        try:
            if progress_cb:
                progress_cb(50, "未检测到办公软件，使用内置引擎转换（排版可能简化）...")
            ok = render_docx_to_pdf(inp, out, progress_cb)
            if ok:
                return True, "内置引擎"
            msgs.append("内置引擎渲染失败")
        except Exception as e:
            msgs.append(f"内置引擎不可用({type(e).__name__}: {e})")
    else:
        msgs.append(".doc 旧格式需本机安装 Office/WPS/LibreOffice")
    raise RuntimeError("；".join(msgs))


def ppt_to_pdf(inp, out, progress_cb=None):
    """PPT/PPTX/DPS → PDF 多级降级。返回 (ok, message)。"""
    ext = os.path.splitext(inp)[1].lower()
    msgs = []
    # 1) Microsoft PowerPoint COM
    try:
        app, progid = _com_dispatch(["PowerPoint.Application"])
        if app is not None:
            if progress_cb:
                progress_cb(40, "通过 PowerPoint 导出 PDF...")
            ok = _ppt_com(app, inp, out)
            _com_cleanup(app)
            if ok and os.path.isfile(out):
                return True, "PowerPoint"
            msgs.append("PowerPoint COM 导出失败")
    except Exception as e:
        msgs.append(f"PowerPoint 不可用({type(e).__name__})")
    # 2) WPS 演示 COM
    try:
        app, progid = _com_dispatch(["Kwpp.Application", "wpp.Application"])
        if app is not None:
            if progress_cb:
                progress_cb(40, "通过 WPS 演示导出 PDF...")
            ok = _ppt_com(app, inp, out)
            _com_cleanup(app)
            if ok and os.path.isfile(out):
                return True, "WPS演示"
            msgs.append("WPS 演示 COM 导出失败")
    except Exception as e:
        msgs.append(f"WPS 演示不可用({type(e).__name__})")
    # 3) LibreOffice headless
    if _soffice_available():
        if progress_cb:
            progress_cb(40, "通过 LibreOffice 导出 PDF...")
        if _libreoffice_convert(inp, out):
            return True, "LibreOffice"
        msgs.append("LibreOffice 转换失败")
    else:
        msgs.append("未安装 LibreOffice")
    # 4) 纯 Python（仅 pptx 可读；老版 .ppt 二进制 python-pptx 无法读取）
    if ext in (".pptx", ".dps"):
        try:
            if progress_cb:
                progress_cb(50, "未检测到办公软件，使用内置引擎转换（排版可能简化）...")
            ok = render_pptx_to_pdf(inp, out, progress_cb)
            if ok:
                return True, "内置引擎"
            msgs.append("内置引擎渲染失败")
        except Exception as e:
            msgs.append(f"内置引擎不可用({type(e).__name__}: {e})")
    else:
        msgs.append(".ppt 旧格式需本机安装 Office/WPS/LibreOffice")
    raise RuntimeError("；".join(msgs))


def render_xlsx_to_pdf(inp, out, progress_cb=None):
    """纯 Python：xlsx → PDF（openpyxl 读表 + reportlab 表格渲染，每表一页）。

    不依赖 Office/LibreOffice。仅支持 .xlsx/.xlsm（openpyxl 可读）；老版 .xls
    二进制格式需 LibreOffice。复杂样式（合并单元格/公式值/图表）会简化。
    """
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
                                    Paragraph, Spacer)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    font_name, _ = _register_cjk_font()
    if progress_cb:
        progress_cb(20, "读取表格...")

    import openpyxl
    wb = openpyxl.load_workbook(inp, read_only=True, data_only=True)
    pdf = SimpleDocTemplate(out, pagesize=landscape(A4),
                            leftMargin=12 * mm, rightMargin=12 * mm,
                            topMargin=12 * mm, bottomMargin=12 * mm)
    styles = getSampleStyleSheet()
    cell_style = ParagraphStyle("Cell", parent=styles["Normal"],
                                fontName=font_name, fontSize=8, leading=10)
    title_style = ParagraphStyle("Sheet", parent=styles["Heading2"],
                                 fontName=font_name, fontSize=13, spaceAfter=6)
    story = []

    def _safe(v):
        s = "" if v is None else str(v)
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    sheet_names = wb.sheetnames
    for si, name in enumerate(sheet_names):
        ws = wb[name]
        story.append(Paragraph(_safe(name), title_style))
        rows = []
        max_col = 0
        for row in ws.iter_rows(values_only=True):
            rows.append([_safe(c) for c in row])
            max_col = max(max_col, len(row))
        if not rows:
            story.append(Spacer(1, 6))
            continue
        # 限制列宽避免超页：每页最多 14 列
        if max_col > 14:
            rows = [[c[:24] for c in r[:14]] for r in rows]
        table = Table(rows, repeatRows=1)
        table.setStyle(TableStyle([
            ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
            ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.93, 0.93, 0.93)),
            ("FONTNAME", (0, 0), (-1, -1), font_name),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3),
            ("RIGHTPADDING", (0, 0), (-1, -1), 3),
        ]))
        story.append(table)
        story.append(Spacer(1, 14))
        if progress_cb:
            pct = 30 + int(si * 55 / max(len(sheet_names), 1))
            progress_cb(min(pct, 88), f"表格 {si + 1}/{len(sheet_names)}...")
    if progress_cb:
        progress_cb(90, "生成PDF...")
    pdf.build(story)
    return os.path.isfile(out)


def excel_to_pdf(inp, out, progress_cb=None):
    """XLS/XLSX → PDF 多级降级。返回 (ok, message)。

    Excel COM 无稳定跨版本 API，优先 LibreOffice；无办公软件时回退 openpyxl
    纯 Python 渲染（排版简化）。
    """
    msgs = []
    ext = os.path.splitext(inp)[1].lower()
    # 1) LibreOffice headless
    if _soffice_available():
        if progress_cb:
            progress_cb(40, "通过 LibreOffice 导出 PDF...")
        if _libreoffice_convert(inp, out):
            return True, "LibreOffice"
        msgs.append("LibreOffice 转换失败")
    else:
        msgs.append("未安装 LibreOffice")
    # 2) 纯 Python（仅 xlsx/xlsm 可读）
    if ext in (".xlsx", ".xlsm"):
        try:
            if progress_cb:
                progress_cb(50, "使用内置引擎转换（排版可能简化）...")
            ok = render_xlsx_to_pdf(inp, out, progress_cb)
            if ok:
                return True, "内置引擎"
            msgs.append("内置引擎渲染失败")
        except Exception as e:
            msgs.append(f"内置引擎不可用({type(e).__name__}: {e})")
    else:
        msgs.append(".xls 旧格式需本机安装 LibreOffice")
    raise RuntimeError("；".join(msgs))
