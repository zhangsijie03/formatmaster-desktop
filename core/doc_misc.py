"""Misc 格式转换 (PPT/HTML/MD/EPUB/RTF/ODT/OFD/图片/TXT/交叉)"""
import os
import re
import shutil
import subprocess
import sys


class DocMiscMixin:
    """PPT/HTML/MD/EPUB/RTF/ODT/OFD/图片/TXT/交叉格式转换方法"""

    # ========== TXT 转换 ==========

    def _txt_to_xlsx(self, inp, out, cb):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        for i, line in enumerate(self._read_text(inp).splitlines()):
            if self._cancel: return False
            parts = line.strip().split('\t')
            if len(parts) == 1:
                parts = line.strip().split(',')
            ws.append(parts)
            if cb and i % 100 == 0:
                cb(20 + min(70, i // 10), f"写入第{i+1}行...")
        wb.save(out)
        return True

    def _txt_to_docx(self, inp, out, cb):
        if cb: cb(30, "生成Word...")
        import docx
        doc = docx.Document()
        for line in self._read_text(inp).splitlines():
            if self._cancel: return False
            text = line.rstrip()
            if text.startswith('# '):
                doc.add_heading(text[2:], level=1)
            elif text.startswith('## '):
                doc.add_heading(text[3:], level=2)
            elif text.startswith('### '):
                doc.add_heading(text[4:], level=3)
            elif text.strip() == '':
                doc.add_paragraph('')
            else:
                doc.add_paragraph(text)
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _txt_to_pptx(self, inp, out, cb):
        if cb: cb(20, "生成PPT...")
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        lines = [line.rstrip() for line in self._read_text(inp).splitlines()
                 if line.strip()]
        chunk_size = 12
        for i in range(0, max(len(lines), 1), chunk_size):
            if self._cancel: return False
            chunk = lines[i:i+chunk_size]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, line in enumerate(chunk):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                if line.startswith('# '):
                    p.text = line[2:]
                    p.font.size = Pt(28)
                    p.font.bold = True
                else:
                    p.text = line[:200]
                    p.font.size = Pt(18)
            if cb:
                cb(20 + int(i * 70 / max(len(lines), 1)), f"生成幻灯片...")
        prs.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _txt_to_pdf(self, inp, out, cb):
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        # 复用 Office PDF 引擎的跨平台 CJK 字体降级链，避免
        # 系统不存在指定 TTF 时中文被 Helvetica 渲染为方框。
        from core.doc_office_pdf import _register_cjk_font
        font_name, _ = _register_cjk_font()
        if cb: cb(30, "读取文本...")
        lines = self._read_text(inp).splitlines()
        pdf_doc = SimpleDocTemplate(out, pagesize=A4)
        styles = getSampleStyleSheet()
        style = styles['Normal']
        style.fontName = font_name
        story = []
        for line in lines:
            safe = line.rstrip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            if safe:
                story.append(Paragraph(safe, style))
            else:
                story.append(Spacer(1, 12))
        if cb: cb(70, "生成PDF...")
        pdf_doc.build(story)
        return True

    def _txt_to_html(self, inp, out, cb):
        if cb: cb(20, "转换中...")
        lines = self._read_text(inp).splitlines()
        body_parts = []
        for line in lines:
            text = line.rstrip()
            if not text:
                body_parts.append('<br>')
            elif text.startswith('# '):
                body_parts.append(f'<h1>{self._safe_html(text[2:])}</h1>')
            elif text.startswith('## '):
                body_parts.append(f'<h2>{self._safe_html(text[3:])}</h2>')
            elif text.startswith('### '):
                body_parts.append(f'<h3>{self._safe_html(text[4:])}</h3>')
            else:
                body_parts.append(f'<p>{self._safe_html(text)}</p>')
        self._write_text(out, self._build_html_page('\n'.join(body_parts)))
        if cb: cb(100, "转换完成")
        return True

    def _txt_to_md(self, inp, out, cb):
        if cb: cb(30, "转换中...")
        shutil.copy2(inp, out)
        if cb: cb(100, "转换完成")
        return True

    # ========== PPT 转换 ==========

    def _pptx_to_pdf(self, inp, out, cb):
        """PPTX/DPS → PDF（多级降级：PowerPoint COM → WPS演示 → LibreOffice → 内置引擎）。

        优先本机办公软件保真排版；无办公软件时自动降级内置引擎渲染。
        """
        if cb: cb(30, "准备导出PDF...")
        from core.doc_office_pdf import ppt_to_pdf
        try:
            ok, engine = ppt_to_pdf(inp, out, cb)
            if cb:
                cb(100, f"转换完成（{engine}）")
            return ok
        except Exception as e:
            if cb:
                cb(-1, f"错误: {str(e)[:200]}")
            return False

    def _pptx_to_txt(self, inp, out, cb):
        from pptx import Presentation
        prs = Presentation(inp)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)
        self._write_text(out, '\n\n'.join(text_parts))
        if cb: cb(100, "转换完成")
        return True

    def _pptx_to_image(self, inp, out, cb):
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        prs = Presentation(inp)
        out_dir = os.path.dirname(out)
        base = os.path.splitext(os.path.basename(out))[0]
        ext = os.path.splitext(out)[1]
        total = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            img = Image.new('RGB', (1280, 720), 'white')
            draw = ImageDraw.Draw(img)
            y = 50
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    for line in shape.text.split('\n'):
                        if line.strip():
                            try:
                                font = ImageFont.truetype("msyh.ttc", 18)
                            except Exception:
                                font = ImageFont.load_default()
                            draw.text((60, y), line.strip(), fill='black', font=font)
                            y += 30
                            if y > 680: break
            if total == 1:
                img.save(out)
            else:
                img.save(os.path.join(out_dir, f"{base}_{i+1}{ext}"))
            if cb:
                cb(20 + int(i * 70 / total), f"幻灯片{i+1}/{total}...")
        return True

    def _pptx_to_docx(self, inp, out, cb):
        if cb: cb(30, "读取演示文稿...")
        from pptx import Presentation
        import docx
        prs = Presentation(inp)
        doc = docx.Document()
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            doc.add_heading(f'幻灯片 {i+1}', level=2)
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    doc.add_paragraph(shape.text.strip())
            if cb:
                cb(20 + int(i * 70 / max(len(prs.slides), 1)), f"幻灯片{i+1}...")
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _pptx_to_html(self, inp, out, cb):
        if cb: cb(30, "读取演示文稿...")
        from pptx import Presentation
        prs = Presentation(inp)
        body_parts = []
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            body_parts.append(f'<h2>幻灯片 {i+1}</h2>')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    body_parts.append(f'<p>{self._safe_html(shape.text.strip())}</p>')
            if cb:
                cb(20 + int(i * 70 / max(len(prs.slides), 1)), f"幻灯片{i+1}...")
        self._write_text(out, self._build_html_page('\n'.join(body_parts)))
        return True

    def _pptx_to_md(self, inp, out, cb):
        if cb: cb(30, "读取演示文稿...")
        from pptx import Presentation
        prs = Presentation(inp)
        lines = []
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            lines.append(f'## 幻灯片 {i+1}')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    lines.append(shape.text.strip())
            lines.append('')
            if cb:
                cb(20 + int(i * 70 / max(len(prs.slides), 1)), f"幻灯片{i+1}...")
        self._write_text(out, '\n'.join(lines))
        return True

    def _ppt_to_pdf(self, inp, out, cb):
        """PPT → PDF（多级降级：PowerPoint COM → WPS演示 → LibreOffice）。

        .ppt 为旧版二进制格式，内置引擎（python-pptx）无法读取；
        必须依赖本机 Office/WPS/LibreOffice，全部缺失时给出明确提示。
        """
        if cb: cb(30, "准备导出PDF...")
        from core.doc_office_pdf import ppt_to_pdf
        try:
            ok, engine = ppt_to_pdf(inp, out, cb)
            if cb:
                cb(100, f"转换完成（{engine}）")
            return ok
        except Exception as e:
            if cb:
                cb(-1, f"错误: {str(e)[:200]}")
            return False

    def _ppt_to_txt(self, inp, out, cb):
        if cb: cb(30, "启动PowerPoint...")
        if sys.platform != "win32":
            # macOS/Linux 没有 PowerPoint COM；优先复用已有的 LibreOffice
            # headless 能力导出纯文本，避免此格式在非 Windows 上直接导入
            # win32com 后失败。
            from core.doc_office_pdf import _find_soffice
            soffice = _find_soffice()
            if not soffice:
                if cb:
                    cb(-1, "旧版 PPT 转 TXT 需要安装 LibreOffice")
                return False
            out_dir = os.path.dirname(os.path.abspath(out)) or "."
            expected = os.path.join(
                out_dir, os.path.splitext(os.path.basename(inp))[0] + ".txt")
            try:
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "txt:Text",
                     "--outdir", out_dir, os.path.abspath(inp)],
                    capture_output=True, timeout=120,
                    creationflags=(0x08000000
                                   if sys.platform == "win32" else 0),
                    encoding="utf-8", errors="ignore")
                if result.returncode != 0 or not os.path.isfile(expected):
                    if cb:
                        cb(-1, "LibreOffice 无法读取旧版 PPT")
                    return False
                if os.path.abspath(expected) != os.path.abspath(out):
                    shutil.move(expected, out)
                if cb: cb(100, "转换完成（LibreOffice）")
                return True
            except Exception as exc:  # noqa: BLE001
                if cb:
                    cb(-1, f"LibreOffice 转换失败：{str(exc)[:160]}")
                return False
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        ppt = None
        pres = None
        try:
            ppt = win32com.client.DispatchEx("PowerPoint.Application")
            ppt.Visible = False
            if cb: cb(50, "读取演示文稿...")
            abs_inp = os.path.abspath(inp)
            pres = ppt.Presentations.Open(abs_inp, WithWindow=False)
            texts = []
            for slide in pres.Slides:
                if self._cancel: return False
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        texts.append(shape.TextFrame.TextRange.Text)
            pres.Close()
            pres = None
            ppt.Quit()
            ppt = None
            self._write_text(out, '\n\n'.join(texts))
            if cb: cb(100, "转换完成")
            return True
        except Exception as e:
            if pres:
                try: pres.Close()
                except Exception: pass
            if ppt:
                try: ppt.Quit()
                except Exception: pass
            raise
        finally:
            pythoncom.CoUninitialize()

    # ========== 图片 转换 ==========

    def _image_to_pdf(self, inp, out, cb):
        """将图片按原始像素无损嵌入 PDF，不做二次有损压缩。

        Pillow 的 ``Image.save(..., "PDF")`` 会把 RGB PNG 重新编码为 JPEG，
        表格、截图和小字号文字会出现明显压缩噪点。PyMuPDF 可直接保留 PNG
        的无损像素流，同时让 JPEG 保持原始编码，避免无意义的质量损失。
        """
        import math
        import pymupdf
        from PIL import Image
        if cb: cb(30, "处理图片...")
        with Image.open(inp) as img:
            width_px, height_px = img.size
            dpi = img.info.get("dpi") or (96.0, 96.0)

        if width_px <= 0 or height_px <= 0:
            raise ValueError("图片尺寸无效")
        if not isinstance(dpi, (tuple, list)):
            dpi = (dpi, dpi)
        dpi_x = float(dpi[0]) if len(dpi) > 0 else 96.0
        dpi_y = float(dpi[1]) if len(dpi) > 1 else dpi_x
        # 外部图片可能携带 0、NaN 或异常 DPI；此时采用屏幕图像通用的
        # 96 DPI，防止生成尺寸无限大或无法打开的 PDF 页面。
        if not math.isfinite(dpi_x) or not 36 <= dpi_x <= 1200:
            dpi_x = 96.0
        if not math.isfinite(dpi_y) or not 36 <= dpi_y <= 1200:
            dpi_y = 96.0

        page_width = width_px * 72.0 / dpi_x
        page_height = height_px * 72.0 / dpi_y
        # PDF 页面边长上限按 200 英寸控制；超大图片只缩小页面物理尺寸，
        # 原始像素仍完整嵌入，不进行重采样。
        max_page_points = 14_400.0
        page_scale = min(1.0, max_page_points / max(page_width, page_height))
        page_width *= page_scale
        page_height *= page_scale

        if cb: cb(70, "保存PDF...")
        pdf = pymupdf.open()
        try:
            page = pdf.new_page(width=page_width, height=page_height)
            page.insert_image(page.rect, filename=inp, keep_proportion=False)
            pdf.set_metadata({"title": os.path.splitext(os.path.basename(inp))[0]})
            pdf.save(out, deflate=True, garbage=4)
        finally:
            pdf.close()
        return True

    def _image_to_docx(self, inp, out, cb):
        if cb: cb(30, "生成Word...")
        import docx
        doc = docx.Document()
        doc.add_picture(inp)
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    # ========== HTML 转换 ==========

    def _html_to_pdf(self, inp, out, cb):
        if cb: cb(30, "解析HTML...")
        content = self._read_text(inp)
        text = re.sub(r'<[^>]+>', ' ', content)
        text = re.sub(r'\s+', ' ', text).strip()
        temp = out + '.tmp.txt'
        self._write_text(temp, text)
        result = self._txt_to_pdf(temp, out, cb)
        try:
            os.remove(temp)
        except OSError:
            pass
        return result

    def _html_to_docx(self, inp, out, cb):
        if cb: cb(30, "解析HTML...")
        import docx
        content = self._read_text(inp)
        doc = docx.Document()
        for match in re.finditer(r'<h([1-3])[^>]*>(.*?)</h\1>', content, re.DOTALL):
            doc.add_heading(re.sub(r'<[^>]+>', '', match.group(2)).strip(), level=int(match.group(1)))
        for match in re.finditer(r'<p[^>]*>(.*?)</p>', content, re.DOTALL):
            text = re.sub(r'<[^>]+>', '', match.group(1)).strip()
            if text:
                doc.add_paragraph(text)
        if not doc.paragraphs:
            text = re.sub(r'<[^>]+>', '', content)
            for line in text.strip().split('\n'):
                if line.strip():
                    doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _html_to_txt(self, inp, out, cb):
        if cb: cb(30, "解析HTML...")
        content = self._read_text(inp)
        text = re.sub(r'<[^>]+>', '', content)
        text = re.sub(r'\s+', ' ', text).strip()
        self._write_text(out, text)
        if cb: cb(100, "转换完成")
        return True

    def _html_to_md(self, inp, out, cb):
        if cb: cb(30, "转换中...")
        content = self._read_text(inp)
        text = re.sub(r'<[^>]+>', '', content)
        text = re.sub(r'\s+', ' ', text).strip()
        self._write_text(out, text)
        if cb: cb(100, "转换完成")
        return True

    def _html_to_xlsx(self, inp, out, cb):
        if cb: cb(30, "解析HTML...")
        import openpyxl
        content = self._read_text(inp)
        wb = openpyxl.Workbook()
        ws = wb.active
        tables = re.findall(r'<table[^>]*>(.*?)</table>', content, re.DOTALL)
        if tables:
            for table_html in tables:
                rows = re.findall(r'<tr[^>]*>(.*?)</tr>', table_html, re.DOTALL)
                for ri, row_html in enumerate(rows):
                    cells = re.findall(r'<t[dh][^>]*>(.*?)</t[dh]>', row_html, re.DOTALL)
                    for ci, cell in enumerate(cells):
                        ws.cell(row=ri+1, column=ci+1, value=re.sub(r'<[^>]+>', '', cell).strip())
                break
        else:
            text = re.sub(r'<[^>]+>', '', content)
            for i, line in enumerate(text.strip().split('\n')[:100]):
                if line.strip():
                    ws.cell(row=i+1, column=1, value=line.strip())
        wb.save(out)
        if cb: cb(100, "转换完成")
        return True

    # ========== Markdown 转换 ==========

    def _md_to_html(self, inp, out, cb):
        if cb: cb(20, "解析Markdown...")
        from markdown_it import MarkdownIt
        md = MarkdownIt()
        html = md.render(self._read_text(inp))
        self._write_text(out, self._build_html_page(html))
        if cb: cb(100, "转换完成")
        return True

    def _md_to_pdf(self, inp, out, cb):
        if cb: cb(20, "解析Markdown...")
        from markdown_it import MarkdownIt
        md = MarkdownIt()
        html = md.render(self._read_text(inp))
        temp_html = out + '.tmp.html'
        self._write_text(temp_html, self._build_html_page(html))
        text = re.sub(r'<[^>]+>', ' ', html)
        text = re.sub(r'\s+', ' ', text).strip()
        temp_txt = out + '.tmp.txt'
        self._write_text(temp_txt, text)
        result = self._txt_to_pdf(temp_txt, out, cb)
        try:
            os.remove(temp_html)
            os.remove(temp_txt)
        except OSError:
            pass
        return result

    def _md_to_docx(self, inp, out, cb):
        if cb: cb(20, "解析Markdown...")
        from markdown_it import MarkdownIt
        import docx
        md = MarkdownIt()
        doc = docx.Document()
        html = md.render(self._read_text(inp))
        for match in re.finditer(r'<h([1-3])[^>]*>(.*?)</h\1>', html, re.DOTALL):
            doc.add_heading(re.sub(r'<[^>]+>', '', match.group(2)).strip(), level=int(match.group(1)))
        for match in re.finditer(r'<li>(.*?)</li>', html, re.DOTALL):
            text = re.sub(r'<[^>]+>', '', match.group(1)).strip()
            if text:
                doc.add_paragraph(text, style='List Bullet')
        for match in re.finditer(r'<p>(.*?)</p>', html, re.DOTALL):
            text = re.sub(r'<[^>]+>', '', match.group(1)).strip()
            if text and not any(doc.paragraphs[-1].text == text for _ in [0] if doc.paragraphs):
                doc.add_paragraph(text)
        if not doc.paragraphs:
            text = re.sub(r'<[^>]+>', '', html)
            for line in text.strip().split('\n'):
                if line.strip():
                    doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _md_to_txt(self, inp, out, cb):
        if cb: cb(20, "解析Markdown...")
        from markdown_it import MarkdownIt
        md = MarkdownIt()
        html = md.render(self._read_text(inp))
        text = re.sub(r'<[^>]+>', '', html)
        text = re.sub(r'\s+', ' ', text).strip()
        self._write_text(out, text)
        if cb: cb(100, "转换完成")
        return True

    # ========== EPUB 转换 ==========

    def _epub_to_text(self, inp):
        text_chunks = []
        for name, raw in self._read_epub(inp):
            try:
                from lxml import html as lx_html
                doc = lx_html.fromstring(raw)
                body = doc.xpath("//body")
                if body:
                    lines = [t.strip() for t in body[0].itertext() if t.strip()]
                    text_chunks.append("\n".join(lines))
            except Exception:  # noqa: BLE001
                continue
        return "\n\n".join(text_chunks)

    def _epub_to_pdf(self, inp, out, cb):
        if cb: cb(20, "读取EPUB...")
        text = self._epub_to_text(inp)
        temp = out + '.tmp.txt'
        self._write_text(temp, text)
        result = self._txt_to_pdf(temp, out, cb)
        try: os.remove(temp)
        except Exception: pass
        return result

    def _epub_to_txt(self, inp, out, cb):
        if cb: cb(20, "读取EPUB...")
        text = self._epub_to_text(inp)
        self._write_text(out, text)
        if cb: cb(100, "转换完成")
        return True

    def _epub_to_html(self, inp, out, cb):
        if cb: cb(20, "读取EPUB...")
        bodies = []
        for name, raw in self._read_epub(inp):
            bodies.append(raw.decode("utf-8", errors="replace"))
        self._write_text(out, self._build_html_page("\n".join(bodies)))
        if cb: cb(100, "转换完成")
        return True

    def _epub_to_docx(self, inp, out, cb):
        if cb: cb(20, "读取EPUB...")
        import docx
        text = self._epub_to_text(inp)
        doc = docx.Document()
        for line in text.split('\n'):
            if line.strip():
                doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    # ========== RTF 转换 ==========

    def _rtf_read(self, inp):
        from striprtf.striprtf import rtf_to_text
        text = self._read_text(inp)
        return rtf_to_text(text)

    def _rtf_to_txt(self, inp, out, cb):
        if cb: cb(30, "解析RTF...")
        self._write_text(out, self._rtf_read(inp))
        if cb: cb(100, "转换完成")
        return True

    def _rtf_to_pdf(self, inp, out, cb):
        if cb: cb(30, "解析RTF...")
        text = self._rtf_read(inp)
        temp = out + '.tmp.txt'
        self._write_text(temp, text)
        result = self._txt_to_pdf(temp, out, cb)
        try: os.remove(temp)
        except Exception: pass
        return result

    def _rtf_to_docx(self, inp, out, cb):
        if cb: cb(30, "解析RTF...")
        import docx
        text = self._rtf_read(inp)
        doc = docx.Document()
        for line in text.split('\n'):
            if line.strip():
                doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    # ========== ODT 转换 ==========

    def _odt_read_text(self, inp):
        from odf.opendocument import load
        from odf.text import P
        doc = load(inp)
        texts = []
        for p in doc.getElementsByType(P):
            texts.append(str(p))
        return '\n'.join(texts)

    def _odt_to_pdf(self, inp, out, cb):
        if cb: cb(30, "读取ODT...")
        text = self._odt_read_text(inp)
        temp = out + '.tmp.txt'
        self._write_text(temp, text)
        result = self._txt_to_pdf(temp, out, cb)
        try: os.remove(temp)
        except Exception: pass
        return result

    def _odt_to_docx(self, inp, out, cb):
        if cb: cb(30, "读取ODT...")
        import docx
        text = self._odt_read_text(inp)
        doc = docx.Document()
        for line in text.split('\n'):
            if line.strip():
                doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _odt_to_txt(self, inp, out, cb):
        if cb: cb(30, "读取ODT...")
        self._write_text(out, self._odt_read_text(inp))
        if cb: cb(100, "转换完成")
        return True

    # ── OFD → PDF ────────────────────────────────────────

    def _ofd_to_pdf(self, inp, out, cb):
        if cb: cb(20, "解析OFD...")
        import pymupdf
        doc = pymupdf.open(inp)
        if cb: cb(50, "转换中...")
        doc.save(out, deflate=True, garbage=4)
        doc.close()
        if cb: cb(100, "转换完成")
        return True

    # ── Word ↔ Excel ─────────────────────────────────────

    def _docx_to_xlsx(self, inp, out, cb):
        if cb: cb(20, "读取Word...")
        import docx
        from openpyxl import Workbook
        doc = docx.Document(inp)
        wb = Workbook()
        if not doc.tables:
            ws = wb.active
            ws.title = "内容"
            for para in doc.paragraphs:
                if para.text.strip():
                    ws.append([para.text.strip()])
        else:
            for ti, table in enumerate(doc.tables):
                ws = wb.create_sheet(title=f"表格{ti+1}") if ti > 0 else wb.active
                ws.title = f"表格{ti+1}"
                for row in table.rows:
                    ws.append([cell.text for cell in row.cells])
        if not wb.sheetnames:
            wb.active.title = "内容"
        wb.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _xlsx_to_docx(self, inp, out, cb):
        if cb: cb(20, "读取Excel...")
        from openpyxl import load_workbook
        import docx
        wb = load_workbook(inp, read_only=True, data_only=True)
        doc = docx.Document()
        for si, sname in enumerate(wb.sheetnames):
            if si > 0:
                doc.add_page_break()
            doc.add_heading(sname, level=2)
            ws = wb[sname]
            rows = list(ws.iter_rows(values_only=True))
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]) if rows[0] else 1)
                table.style = "Table Grid"
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        if ci < len(table.columns):
                            table.cell(ri, ci).text = str(val) if val is not None else ""
        wb.close()
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True
